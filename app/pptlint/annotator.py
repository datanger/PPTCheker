"""
带标记PPT输出（对应任务：实现PPT注释输出模块并集成CLI）

实现要点：
- 在每页左上角新增“问题汇总”文本框；
- 对命中的 shape，将其文本末尾追加“【标记: 规则ID】”；
- 不覆盖原文件，另存为副本。
"""
from collections import defaultdict
from typing import List, Optional
from pptx import Presentation
from pptx.util import Pt, Inches

from .model import Issue
from .llm import LLMClient


def _contains_acronym(text: str) -> bool:
    """检查文本是否包含需要解释的缩略语（已废弃，保留用于向后兼容）"""
    # 注意：此函数已被废弃，缩略语识别现在完全由LLM进行
    # 保留此函数仅用于向后兼容，实际不再使用
    return False


def _is_acronym_adequately_explained(text: str, acronym: str, llm_client: Optional[LLMClient] = None) -> bool:
    """使用LLM判断缩略语是否已经被充分解释"""
    if llm_client is None:
        # 如果没有LLM客户端，使用改进的启发式判断
        # 检查是否包含冒号、括号等解释性标点
        explanation_indicators = [':', '：', '(', '（', '（', '）', '是', '为', '指', '即']
        
        # 检查是否有明确的解释模式
        # 模式1：缩略语：全称
        if f"{acronym}：" in text or f"{acronym}:" in text:
            return True
        
        # 模式2：缩略语（全称）
        if f"{acronym}（" in text or f"{acronym}(" in text:
            return True
        
        # 模式3：全称（缩略语）
        if f"（{acronym}）" in text or f"({acronym})" in text:
            return True
        
        # 模式4：包含解释性词汇
        if any(indicator in text for indicator in explanation_indicators):
            # 进一步检查是否在缩略语附近有解释
            import re
            # 查找缩略语附近的文本（前后20个字符）
            pattern = rf".{{0,20}}{acronym}.{{0,20}}"
            matches = re.findall(pattern, text)
            for match in matches:
                if any(indicator in match for indicator in explanation_indicators):
                    return True
        
        return False
    
    try:
        # 构建LLM提示词
        prompt = f"""请判断以下文本中的缩略语"{acronym}"是否已经被充分解释。

文本内容：
{text}

请分析：
1. 该缩略语是否出现
2. 是否提供了完整的解释（包括全称和含义）
3. 解释是否清晰易懂

请只回答"是"或"否"。

回答："""
        
        response = llm_client.chat(prompt)
        # 清理响应，提取"是"或"否"
        response_text = response.strip().lower()
        if '是' in response_text and '否' not in response_text:
            return True
        elif '否' in response_text and '是' not in response_text:
            return False
        else:
            # 如果LLM回答不明确，使用改进的启发式判断
            explanation_indicators = [':', '：', '(', '（', '（', '）', '是', '为', '指', '即']
            
            # 检查是否有明确的解释模式
            if f"{acronym}：" in text or f"{acronym}:" in text:
                return True
            
            if f"{acronym}（" in text or f"{acronym}(" in text:
                return True
            
            if f"（{acronym}）" in text or f"({acronym})" in text:
                return True
            
            return any(indicator in text for indicator in explanation_indicators)
            
    except Exception as e:
        print(f"LLM判断缩略语解释失败: {e}")
        # 回退到改进的启发式判断
        explanation_indicators = [':', '：', '(', '（', '（', '）', '是', '为', '指', '即']
        
        # 检查是否有明确的解释模式
        if f"{acronym}：" in text or f"{acronym}:" in text:
            return True
        
        if f"{acronym}（" in text or f"{acronym}(" in text:
            return True
        
        if f"（{acronym}）" in text or f"({acronym})" in text:
            return True
        
        return any(indicator in text for indicator in explanation_indicators)


def annotate_pptx(src_path: str, issues: List[Issue], output_path: str, llm_client: Optional[LLMClient] = None) -> None:
    prs = Presentation(src_path)

    # 按页聚合问题
    issues_by_slide = defaultdict(list)
    for it in issues:
        issues_by_slide[it.slide_index].append(it)

    # 全局问题汇总：包含所有问题类型，不过滤info级别
    from collections import Counter
    rule_to_label = {
        # 规则检查问题
        "FontFamilyRule": "字体不规范",
        "FontSizeRule": "字号过小",
        "ColorCountRule": "颜色过多",
        "ThemeHarmonyRule": "色调不一致",
        # LLM智能审查问题
        "LLM_AcronymRule": "专业缩略语需解释",
        "LLM_ContentRule": "内容逻辑问题",
        "LLM_FormatRule": "智能格式问题",
        "LLM_FluencyRule": "表达流畅性问题",
        "LLM_TitleStructureRule": "标题结构问题",
    }
    
    # 统计所有问题类型
    grouped_all = Counter((rule_to_label.get(it.rule_id, "其他问题"), it.severity) for it in issues)
    global_summary_lines = [
        f"- {label} [{sev}] x{cnt}"
        for (label, sev), cnt in grouped_all.items()
    ]

    for s_idx, slide in enumerate(prs.slides):
        page_issues = issues_by_slide.get(s_idx, [])

        # 仅在首页绘制全局汇总
        if s_idx == 0 and global_summary_lines:
            left, top, width, height = Inches(0.3), Inches(0.2), Inches(6.5), Inches(1.8)
            tf_box = slide.shapes.add_textbox(left, top, width, height)
            tf = tf_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "问题汇总:\n" + "\n".join(global_summary_lines)
            if run.font is not None:
                run.font.size = Pt(12)

        # 对对象内联标记
        page = prs.slides[s_idx]
        
        # 调试信息：显示该页面的所有问题
        if page_issues:
            print(f"\n页面 {s_idx + 1} 的问题:")
            for issue in page_issues:
                print(f"  - {issue.rule_id}: {issue.object_ref} - {issue.message}")
        
        # 先处理LLM问题：收集该页面的所有LLM问题
        llm_issues = [issue for issue in page_issues if issue.rule_id.startswith("LLM_")]
        if llm_issues:
            print(f"    页面 {s_idx} 发现 {len(llm_issues)} 个LLM问题:")
            for issue in llm_issues:
                print(f"      {issue.rule_id}: {issue.object_ref} - {issue.message}")
        
        for shp in page.shapes:
            # 更安全的属性检查
            if not hasattr(shp, "text_frame") or shp.text_frame is None:
                continue
                
            # 改进对象引用匹配：支持多种引用方式
            sid = str(getattr(shp, "shape_id", ""))
            hit_rules = []
            
            # 不再自动标记所有文本对象，而是根据具体问题类型进行精确匹配
            for issue in page_issues:
                # 匹配方式1：直接shape_id匹配
                if issue.object_ref == sid:
                    hit_rules.append(issue.rule_id)
                # 匹配方式2：text_block_X_Y格式匹配（LLM返回的精确格式）
                elif issue.object_ref.startswith("text_block_") and "_" in issue.object_ref:
                    # text_block_2_1 格式：分割后是 ["text", "block", "2", "1"]
                    # 所以页码是 parts[2]，块索引是 parts[3]
                    parts = issue.object_ref.split("_")
                    if len(parts) >= 4 and parts[2] == str(s_idx):
                        print(f"    🔍 检查text_block匹配: {issue.object_ref} -> 页面 {s_idx}")
                        # 对于text_block格式，我们检查文本内容是否包含相关缩略语
                        if (issue.rule_id == "LLM_AcronymRule" or 
                            issue.rule_id.endswith("_AcronymRule")):
                            # 检查文本内容是否包含缩略语
                            text_content = ""
                            try:
                                for para in shp.text_frame.paragraphs:
                                    for run in para.runs:
                                        text_content += run.text + " "
                            except:
                                text_content = ""
                            
                            print(f"    📝 形状 {sid} 文本内容: {text_content[:50]}...")
                            
                            # 智能检测缩略语是否需要解释
                            if _contains_acronym(text_content):
                                # 提取检测到的缩略语
                                import re
                                potential_acronyms = re.findall(r'\b[A-Z]{2,10}\b', text_content)
                                common_words = {'THE', 'AND', 'FOR', 'ARE', 'BUT', 'NOT', 'YOU', 'ALL', 'CAN', 'HER', 'WAS', 'ONE', 'OUR', 'OUT', 'DAY', 'GET', 'HAS', 'HIM', 'HIS', 'HOW', 'MAN', 'NEW', 'NOW', 'OLD', 'SEE', 'TWO', 'WAY', 'WHO', 'BOY', 'DID', 'ITS', 'LET', 'PUT', 'SAY', 'SHE', 'TOO', 'USE'}
                                acronyms = [acronym for acronym in potential_acronyms if acronym not in common_words]
                                
                                # 检查每个缩略语是否已经被充分解释
                                needs_explanation = False
                                for acronym in acronyms:
                                    if not _is_acronym_adequately_explained(text_content, acronym, llm_client):
                                        needs_explanation = True
                                        print(f"    🔍 缩略语 {acronym} 需要解释")
                                        break
                                
                                if needs_explanation:
                                    hit_rules.append(issue.rule_id)
                                    print(f"    ✅ 智能匹配: 形状 {sid} 包含需要解释的缩略语，标记为 {issue.rule_id}")
                                else:
                                    print(f"    ✅ 形状 {sid} 的缩略语已被充分解释，跳过标记")
                            else:
                                pass
                    else:
                        # 对于其他LLM规则，暂时跳过
                        pass
                # 匹配方式3：title_X格式匹配（页面标题）
                elif issue.object_ref.startswith("title_") and issue.object_ref.endswith(f"_{s_idx}"):
                    # 对于标题问题，我们标记该页面的标题对象
                    if shp.is_title and shp.title_level:
                        hit_rules.append(issue.rule_id)
                        print(f"    标题匹配: 形状 {sid} 是标题，标记为 {issue.rule_id}")
                    elif shp == page.shapes[0]:  # 备用方案：假设第一个形状是标题
                        hit_rules.append(issue.rule_id)
                        print(f"    标题备用匹配: 形状 {sid} 是第一个形状，标记为 {issue.rule_id}")
                # 匹配方式4：page_X格式匹配（页面级别问题）
                elif issue.object_ref.startswith("page_") and issue.object_ref.endswith(f"_{s_idx}"):
                    # 对于页面级别问题，我们需要检查文本内容是否包含相关缩略语
                    if (issue.rule_id == "LLM_AcronymRule" or 
                        issue.rule_id.endswith("_AcronymRule")):
                        # 对于页面级别的缩略语问题，检查当前形状是否包含缩略语
                        print(f"    🔍 检查page_X匹配: {issue.object_ref} -> 页面 {s_idx}")
                        
                        # 获取形状的文本内容
                        text_content = ""
                        try:
                            for para in shp.text_frame.paragraphs:
                                for run in para.runs:
                                    text_content += run.text + " "
                        except:
                            text_content = ""
                        
                        # 如果形状包含缩略语，则标记
                        if text_content.strip() and _contains_acronym(text_content):
                            # 提取检测到的缩略语
                            import re
                            potential_acronyms = re.findall(r'\b[A-Z]{2,10}\b', text_content)
                            common_words = {'THE', 'AND', 'FOR', 'ARE', 'BUT', 'NOT', 'YOU', 'ALL', 'CAN', 'HER', 'WAS', 'ONE', 'OUR', 'OUT', 'DAY', 'GET', 'HAS', 'HIM', 'HIS', 'HOW', 'MAN', 'NEW', 'NOW', 'OLD', 'SEE', 'TWO', 'WAY', 'WHO', 'BOY', 'DID', 'ITS', 'LET', 'PUT', 'SAY', 'SHE', 'TOO', 'USE'}
                            acronyms = [acronym for acronym in potential_acronyms if acronym not in common_words]
                            
                            # 检查每个缩略语是否已经被充分解释
                            needs_explanation = False
                            for acronym in acronyms:
                                if not _is_acronym_adequately_explained(text_content, acronym, llm_client):
                                    needs_explanation = True
                                    print(f"    🔍 页面级别缩略语 {acronym} 需要解释")
                                    break
                            
                            if needs_explanation:
                                hit_rules.append(issue.rule_id)
                                print(f"    ✅ 页面级别智能匹配: 形状 {sid} 包含需要解释的缩略语，标记为 {issue.rule_id}")
                            else:
                                print(f"    ✅ 形状 {sid} 的缩略语已被充分解释，跳过标记")
                        else:
                            print(f"    ❌ 形状 {sid} 不包含缩略语，跳过页面级别标记")
                    else:
                        # 对于其他LLM规则，直接添加
                        hit_rules.append(issue.rule_id)
                # 匹配方式5：page级别的问题（向后兼容）
                elif issue.object_ref == "page":
                    # 对于page级别问题，我们标记该页面的所有文本对象
                    hit_rules.append(issue.rule_id)
                # 匹配方式6：全局缩略语问题（当LLM报告页面级别问题时，检查所有页面）
                elif (issue.rule_id in ["LLM_AcronymRule", "ADAS_AcronymRule", "GraphRAG_AcronymRule"] or 
                      issue.rule_id.endswith("_AcronymRule")) and issue.object_ref.startswith("page_"):
                    # 对于LLM报告的页面级别缩略语问题，检查当前形状是否包含相关缩略语                    
                    # 获取形状的文本内容
                    text_content = ""
                    try:
                        for para in shp.text_frame.paragraphs:
                            for run in para.runs:
                                text_content += run.text + " "
                    except:
                        text_content = ""
                    
                    # 如果形状包含缩略语，则标记
                    if text_content.strip() and _contains_acronym(text_content):
                        # 提取检测到的缩略语
                        import re
                        potential_acronyms = re.findall(r'\b[A-Z]{2,10}\b', text_content)
                        common_words = {'THE', 'AND', 'FOR', 'ARE', 'BUT', 'NOT', 'YOU', 'ALL', 'CAN', 'HER', 'WAS', 'ONE', 'OUR', 'OUT', 'DAY', 'GET', 'HAS', 'HIM', 'HIS', 'HOW', 'MAN', 'NEW', 'NOW', 'OLD', 'SEE', 'TWO', 'WAY', 'WHO', 'BOY', 'DID', 'ITS', 'LET', 'PUT', 'SAY', 'SHE', 'TOO', 'USE'}
                        acronyms = [acronym for acronym in potential_acronyms if acronym not in common_words]
                        
                        # 关键修复：只标记包含目标缩略语的形状
                        # 从issue.message中提取目标缩略语名称
                        target_acronym = None
                        if "ADAS" in issue.message:
                            target_acronym = "ADAS"
                        elif "GraphRAG" in issue.message:
                            target_acronym = "GraphRAG"
                        elif "LLM" in issue.message:
                            target_acronym = "LLM"
                        # 可以继续添加其他缩略语
                        
                        if target_acronym and target_acronym in acronyms:
                            # 检查目标缩略语是否已经被充分解释
                            if not _is_acronym_adequately_explained(text_content, target_acronym, llm_client):
                                hit_rules.append(issue.rule_id)
                                print(f"    ✅ 全局缩略语匹配: 形状 {sid} 包含需要解释的缩略语 {target_acronym}，标记为 {issue.rule_id}")
                            else:
                                print(f"    ✅ 形状 {sid} 的缩略语 {target_acronym} 已被充分解释，跳过标记")
                        else:
                            print(f"    ❌ 形状 {sid} 不包含目标缩略语 {target_acronym}，跳过全局缩略语标记")
                    else:
                        print(f"    ❌ 形状 {sid} 不包含缩略语，跳过全局缩略语标记")
            
            if not hit_rules:
                continue
                
            # 规则到中文类别的映射
            rule_to_label = {
                # 规则检查问题
                "FontFamilyRule": "字体不规范",
                "FontSizeRule": "字号过小",
                "ColorCountRule": "颜色过多",
                "ThemeHarmonyRule": "色调不一致",
                # LLM智能审查问题
                "LLM_AcronymRule": "专业缩略语需解释",
                "ADAS_AcronymRule": "专业缩略语需解释",
                "GraphRAG_AcronymRule": "专业缩略语需解释",
                "LLM_ContentRule": "内容逻辑问题",
                "LLM_FormatRule": "智能格式问题",
                "LLM_FluencyRule": "表达流畅性问题",
                "LLM_TitleStructureRule": "标题结构问题",
            }
            
            # 允许多个不同类别；同类多次命中以 xN 展示
            from collections import Counter
            label_counts = Counter(rule_to_label.get(rid, "其他问题") for rid in hit_rules)
            labels = [f"{lab}x{cnt}" if cnt > 1 else lab for lab, cnt in label_counts.items()]
            
            # 调试信息：显示匹配到的规则
            if hit_rules:
                print(f"页面 {s_idx} 形状 {sid} 匹配到规则: {hit_rules}")
                if any(rid.startswith("LLM_") for rid in hit_rules):
                    print(f"    -> 包含LLM规则，将应用样式和标记")
            
            try:
                # 对现有 runs 施加样式：红色 + 下划线（不倾斜）
                for para in shp.text_frame.paragraphs:
                    for r in para.runs:
                        if r.font is not None:
                            # 取消倾斜
                            r.font.italic = False
                            # 优先设置为波浪线，不支持则退化为普通下划线
                            try:
                                from pptx.enum.text import MSO_TEXT_UNDERLINE
                                r.font.underline = MSO_TEXT_UNDERLINE.WAVY_LINE
                            except Exception:
                                r.font.underline = True
                            # 设为红色
                            try:
                                from pptx.dml.color import RGBColor
                                r.font.color.rgb = RGBColor(255, 0, 0)
                            except Exception:
                                pass
                # 同时在最后追加规则摘要（去重后的中文类别），便于溯源
                para_tail = shp.text_frame.paragraphs[-1]
                tail = para_tail.add_run()
                if labels:
                    tail.text = " 【标记: " + "、".join(labels) + "】"
                else:
                    tail.text = " 【标记: 规范问题】"
                
                # 调试信息：显示标记内容
                print(f"    📝 为形状 {sid} 添加标记: '{tail.text}'")
                
                if tail.font is not None:
                    tail.font.size = Pt(10)
                    # 将标记文字设为蓝色
                    try:
                        from pptx.dml.color import RGBColor
                        tail.font.color.rgb = RGBColor(0, 0, 255)
                        print(f"    🎨 设置标记颜色为蓝色")
                    except Exception as e:
                        print(f"    ⚠️ 设置标记颜色失败: {e}")
                else:
                    print(f"    ⚠️ 形状 {sid} 的标记字体对象为空")
                    
                print(f"    ✅ 形状 {sid} 标记完成")
            except Exception as e:
                # 不阻断流程，记录错误
                print(f"标记形状 {sid} 时出错: {e}")
                pass

    prs.save(output_path)

