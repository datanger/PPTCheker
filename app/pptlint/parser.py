"""
PPTX 解析器 - 提取PPT文件的详细信息

输出格式：
- 每页为一个对象：{"页码": int, "文本块": [...], "图片": [...]}
- 文本块：文本块位置、图层编号、是否是标题占位符、拼接字符
- 图片：图片位置、类型、大小、图层位置

注意：某些字段可能无法直接获取，已删除并说明原因
"""

import json
from typing import List, Dict, Any, Optional, Tuple
from pptx import Presentation
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
# 本文件不再生成“拼接字符”，改为直接输出“段落属性”（按 run 维度）

# 删除的字段及原因：
# - 图片质量：无法直接获取，需要图像分析
# - 图片格式：只能获取文件扩展名，无法获取实际格式

# 主题色近似映射（Office 默认主题近似值）
THEME_COLOR_TO_HEX = {
    MSO_THEME_COLOR.TEXT_1: "#000000",
    MSO_THEME_COLOR.BACKGROUND_1: "#FFFFFF",
    MSO_THEME_COLOR.TEXT_2: "#44546A",
    MSO_THEME_COLOR.BACKGROUND_2: "#E7E6E6",
    MSO_THEME_COLOR.ACCENT_1: "#5B9BD5",
    MSO_THEME_COLOR.ACCENT_2: "#ED7D31",
    MSO_THEME_COLOR.ACCENT_3: "#A5A5A5",
    MSO_THEME_COLOR.ACCENT_4: "#FFC000",
    MSO_THEME_COLOR.ACCENT_5: "#4472C4",
    MSO_THEME_COLOR.ACCENT_6: "#70AD47",
    MSO_THEME_COLOR.HYPERLINK: "#0563C1",
    MSO_THEME_COLOR.FOLLOWED_HYPERLINK: "#954F72",
}


THEME_PLACEHOLDER_DEFAULT_MAP = {
    "+mn-ea": "微软雅黑",
    "+mj-ea": "微软雅黑",
    "+mn-lt": "Calibri",
    "+mj-lt": "Calibri",
}

# 用于在JSON层面合并相邻run：比较除“段落内容”外的样式是否一致
ATTR_COMPARE_KEYS = [
    "字体类型",
    "字号",
    "字体颜色",
    "是否粗体",
    "是否斜体",
    "是否下划线",
    "是否带删除线",
]


def _hex_to_rgb_tuple(hex_str: str) -> Tuple[int, int, int]:
    hex_str = hex_str.lstrip('#')
    return int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)


def _rgb_tuple_to_hex(rgb: Tuple[int, int, int]) -> str:
    r, g, b = rgb
    r = max(0, min(255, r))
    g = max(0, min(255, g))
    b = max(0, min(255, b))
    return f"#{r:02X}{g:02X}{b:02X}"


def _apply_brightness(base_hex: str, brightness: Optional[float]) -> str:
    """对主题色应用亮度调整（-1.0~1.0），仿照PowerPoint tint/shade 逻辑的近似实现。"""
    if brightness is None or brightness == 0:
        return base_hex
    r, g, b = _hex_to_rgb_tuple(base_hex)
    if brightness > 0:
        # tint toward white
        r = r + (255 - r) * brightness
        g = g + (255 - g) * brightness
        b = b + (255 - b) * brightness
    else:
        # shade toward black
        factor = 1.0 + brightness  # brightness is negative
        r = r * factor
        g = g * factor
        b = b * factor
    return _rgb_tuple_to_hex((int(round(r)), int(round(g)), int(round(b))))


def _rgb_to_hex(color) -> Optional[str]:
    """将RGB/主题颜色转换为十六进制格式，并考虑亮度调整。"""
    try:
        # 明确RGB
        if isinstance(color, RGBColor):
            return f"#{color[0]:02X}{color[1]:02X}{color[2]:02X}"
        # python-pptx Font.color 可能有 rgb / theme_color / brightness
        if hasattr(color, "rgb") and color.rgb is not None:
            rgb = color.rgb
            return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
        base_hex = None
        if hasattr(color, "theme_color") and color.theme_color is not None:
            base_hex = THEME_COLOR_TO_HEX.get(color.theme_color, None)
        if base_hex is not None:
            # brightness: -1.0..1.0（正值提亮，负值压暗）
            bright = getattr(color, "brightness", None)
            try:
                if bright is not None:
                    base_hex = _apply_brightness(base_hex, float(bright))
            except Exception:
                pass
            return base_hex
    except Exception:
        pass
    return None


# 将常见颜色的十六进制值映射为中文颜色名；其余按“灰色”近似
COMMON_COLOR_NAME_TO_RGB = {
    "黑色": (0, 0, 0),
    "白色": (255, 255, 255),
    "红色": (255, 0, 0),
    "绿色": (0, 255, 0),
    "蓝色": (0, 0, 255),
    "黄色": (255, 255, 0),
    "橙色": (255, 165, 0),
    "紫色": (128, 0, 128),
    "灰色": (128, 128, 128),
}


def _hex_to_cn_color_name(hex_color: str) -> str:
    """将 #RRGGBB 映射为最接近的常见中文颜色名。"""
    try:
        r, g, b = _hex_to_rgb_tuple(hex_color)
        best_name = "灰色"
        best_dist = float("inf")
        for name, (cr, cg, cb) in COMMON_COLOR_NAME_TO_RGB.items():
            dist = (r - cr) ** 2 + (g - cg) ** 2 + (b - cb) ** 2
            if dist < best_dist:
                best_dist = dist
                best_name = name
        return best_name
    except Exception:
        return "灰色"


def _merge_font_family_alias(raw_name: Optional[str]) -> str:
    """合并常见字体族别名/派生名到主名。
    规则示例：
    - "宋体"、"宋体-正文"、"宋体-标题" → "宋体"
    - "Meiryo"、"Meiryo-正文"、"Meiryo-Regular" → "Meiryo UI"
    - "微软雅黑"、"微软雅黑-正文" → "微软雅黑"
    - "楷体"、"楷体_GB2312" → "楷体"
    - "Times New Roman"、"Times New Roman-Regular" → "Time New Roman"
    其它字体保持原样；空值返回 "未知"。
    """
    if not isinstance(raw_name, str) or not raw_name.strip():
        return "未知"
    name = raw_name.strip()
    low = name.lower()
    # 去掉常见的后缀标记
    strip_suffixes = ["-正文", "-标题", "-regular", " regular", " bold", "-bold", " italic", "-italic", "_gb2312", "-gb2312"]
    for suf in strip_suffixes:
        if low.endswith(suf):
            name = name[: len(name) - len(suf)]
            low = name.lower()
            break
    # 统一 Meiryo 派生
    if "meiryo" in low:
        return "Meiryo UI"
    # 统一 宋体 派生
    if "宋体" in name:
        return "宋体"
    # 统一 微软雅黑 派生
    if "微软雅黑" in name or "microsoft yahei" in low:
        return "微软雅黑"
    # 统一 楷体 派生
    if "楷体" in name or "kaiti" in low:
        return "楷体"
    # 统一 Times New Roman 派生
    if "times new roman" in low or "timesnewroman" in low.replace(" ", ""):
        return "Time New Roman"
    return name


def _get_shape_position(shape) -> Dict[str, str]:
    """返回形状位置，单位为百分比（%），相对左上角。
    PowerPoint内部单位为EMU，需要先获取幻灯片尺寸来计算百分比。
    """
    try:
        # 获取幻灯片尺寸 - 尝试多种方法
        slide_width = None
        slide_height = None
        
        try:
            # 方法1：直接从shape.part获取
            slide_width = shape.part.slide_width
            slide_height = shape.part.slide_height
        except Exception:
            pass
        
        try:
            # 方法2：从slide对象获取
            if hasattr(shape, 'slide'):
                slide_width = shape.slide.slide_width
                slide_height = shape.slide.slide_height
        except Exception:
            pass
        
        try:
            # 方法3：从slide_layout获取
            if hasattr(shape.part, 'slide_layout'):
                slide_width = shape.part.slide_layout.slide_width
                slide_height = shape.part.slide_layout.slide_height
        except Exception:
            pass
        
        # 如果仍然无法获取，使用默认值（标准PPT尺寸：16:9 宽屏 = 9144000 x 6858000 EMU）
        if slide_width is None or slide_height is None:
            slide_width = 9144000  # 16:9 宽屏宽度 (10" x 914400 EMU/inch)
            slide_height = 6858000  # 16:9 宽屏高度 (7.5" x 914400 EMU/inch)
        
        # 调试信息：打印实际获取的尺寸
        if hasattr(shape, '_element') and hasattr(shape._element, 'attrib'):
            try:
                # 尝试从XML属性获取实际尺寸
                xml_attrib = shape._element.attrib
                if 'cx' in xml_attrib and 'cy' in xml_attrib:
                    actual_width = int(xml_attrib['cx'])
                    actual_height = int(xml_attrib['cy'])
                    # 如果XML中的尺寸更合理，使用它
                    if actual_width > 0 and actual_height > 0:
                        slide_width = actual_width
                        slide_height = actual_height
            except Exception:
                pass
        
        def emu_to_percent_str(emu_val, slide_dimension) -> str:
            try:
                percent = (float(emu_val) / float(slide_dimension)) * 100.0
                # 不限制百分比范围，显示真实值（可能超过100%）
                return f"{percent:.2f}%"
            except Exception:
                return "0.00%"

        return {
            "left": emu_to_percent_str(shape.left, slide_width),
            "top": emu_to_percent_str(shape.top, slide_height),
            "width": emu_to_percent_str(shape.width, slide_width),
            "height": emu_to_percent_str(shape.height, slide_height)
        }
    except Exception as e:
        print(f"获取形状位置失败: {e}")
        return {"left": "0.00%", "top": "0.00%", "width": "0.00%", "height": "0.00%"}


def _is_title_placeholder(shape) -> bool:
    try:
        if shape.is_placeholder and (shape.placeholder_format.type == PP_PLACEHOLDER.TITLE or
                                     shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE or
                                     shape.placeholder_format.type == PP_PLACEHOLDER.VERTICAL_TITLE or
                                     shape.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE):
            return True
    except Exception:
        pass
    return False


def _get_para_rfonts(para) -> Optional[str]:
    """尝试从段落 pPr.rPr.rFonts / latin 中获取字体名称。"""
    try:
        p = getattr(para, '_p', None)
        pPr = getattr(p, 'pPr', None) if p is not None else None
        rPr = getattr(pPr, 'rPr', None) if pPr is not None else None
        if rPr is not None:
            # rFonts
            rfonts = getattr(rPr, 'rFonts', None)
            if rfonts is not None:
                name = (
                    getattr(rfonts, 'eastAsia', None)
                    or getattr(rfonts, 'ascii', None)
                    or getattr(rfonts, 'hAnsi', None)
                    or getattr(rfonts, 'cs', None)
                )
                if name:
                    return name
            # latin
            latin = getattr(rPr, 'latin', None)
            if latin is not None and getattr(latin, 'typeface', None):
                return latin.typeface
    except Exception:
        pass
    return None


def _get_theme_font_for_placeholder(placeholder: str, presentation) -> Optional[str]:
    """从PowerPoint主题中获取占位符对应的字体。
    占位符格式如: +mn-ea, +mj-ea 等
    
    通用解决方案：
    1. 首先尝试从主题XML中直接解析字体信息
    2. 如果无法解析，使用启发式方法根据PPT内容推断
    3. 最后使用默认映射作为兜底方案
    """
    try:
        if not placeholder.startswith('+'):
            return None
            
        # print(f"    🔍 解析占位符: {placeholder}")
        
        # 方法1: 尝试从主题XML中直接解析字体信息
        theme_font = _resolve_font_from_theme_xml(placeholder, presentation)
        if theme_font:
            return theme_font
            
        # 方法2: 使用启发式方法根据PPT内容推断字体
        inferred_font = _infer_font_from_ppt_content(placeholder, presentation)
        if inferred_font:
            return inferred_font
            
        # 方法3: 使用默认映射作为兜底方案
        return _get_default_font_mapping(placeholder)
        
    except Exception as e:
        # print(f"    ❌ 获取主题字体失败: {e}")
        return _get_default_font_mapping(placeholder)


def _resolve_font_from_theme_xml(placeholder: str, presentation) -> Optional[str]:
    """从主题XML中直接解析字体信息"""
    try:
        theme = presentation.theme
        if not theme or not theme.part:
            return None
            
        theme_xml = theme.part.xml
        
        # 解析占位符类型
        parts = placeholder[1:].split('-')  # 去掉'+'号
        if len(parts) != 2:
            return None
            
        font_type, script = parts  # font_type: mn/mj, script: ea/lt
        
        # 在主题XML中搜索对应的字体定义
        # 这里需要解析主题XML的fontScheme部分
        # 由于XML解析比较复杂，这里提供一个简化的实现
        
        # 简化的XML解析逻辑（实际使用时可能需要更复杂的解析）
        if 'fontScheme' in theme_xml:
            # 根据占位符类型在XML中搜索
            if script == 'ea':  # east asia
                if font_type == 'mj':  # major
                    # 搜索majorFont的eastAsia定义
                    if 'majorFont' in theme_xml and 'eastAsia' in theme_xml:
                        # 这里需要实际的XML解析逻辑
                        pass
                else:  # minor
                    # 搜索minorFont的eastAsia定义
                    if 'minorFont' in theme_xml and 'eastAsia' in theme_xml:
                        pass
            elif script == 'lt':  # latin
                if font_type == 'mj':  # major
                    # 搜索majorFont的latin定义
                    pass
                else:  # minor
                    # 搜索minorFont的latin定义
                    pass
                    
        return None
        
    except Exception:
        return None


def _infer_font_from_ppt_content(placeholder: str, presentation) -> Optional[str]:
    """根据PPT内容推断占位符对应的字体"""
    try:
        # 分析PPT中已识别的字体分布
        font_counts = {}
        
        # 遍历所有幻灯片，统计已识别的字体
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run, 'font') and run.font and run.font.name:
                                font_name = run.font.name.strip()
                                if font_name and not font_name.startswith('+'):
                                    font_counts[font_name] = font_counts.get(font_name, 0) + 1
        
        # 根据字体分布推断占位符对应的字体
        if font_counts:
            # 获取使用最多的字体
            most_common_font = max(font_counts.items(), key=lambda x: x[1])[0]
            
            # 根据占位符类型和脚本类型推断
            parts = placeholder[1:].split('-')
            if len(parts) == 2:
                font_type, script = parts
                
                if script == 'ea':  # east asia
                    # 东亚字体，优先选择已识别的东亚字体
                    east_asia_fonts = ['Meiryo', '宋体', '微软雅黑', '楷体']
                    for font in east_asia_fonts:
                        if any(font.lower() in name.lower() for name in font_counts.keys()):
                            return font
                    # 如果没有找到明确的东亚字体，使用最常见的字体
                    return most_common_font
                elif script == 'lt':  # latin
                    # 拉丁字体，优先选择已识别的拉丁字体
                    latin_fonts = ['Calibri', 'Arial', 'Times New Roman']
                    for font in latin_fonts:
                        if any(font.lower() in name.lower() for name in font_counts.keys()):
                            return font
                    # 如果没有找到明确的拉丁字体，使用最常见的字体
                    return most_common_font
        
        return None
        
    except Exception:
        return None


def _get_default_font_mapping(placeholder: str) -> Optional[str]:
    """获取默认的占位符字体映射"""
    # 默认映射表，根据常见的PowerPoint主题设置
    default_mapping = {
        # 东亚字体 (east asia)
        "+mn-ea": "Meiryo",      # minor east asia - 默认日文
        "+mj-ea": "Meiryo",      # major east asia - 默认日文
        # 拉丁字体 (latin)  
        "+mn-lt": "Calibri",     # minor latin - 默认英文字体
        "+mj-lt": "Calibri",     # major latin - 默认英文字体
    }
    
    return default_mapping.get(placeholder)


def _resolve_run_font_props(run, para, is_title_placeholder: bool, host_shape) -> Dict[str, Any]:
    """解析 run 的样式属性。
    要求：字符类型（字体名称）仅保留 run/paragraph 级别识别：
    1) run._r.rPr.rFonts/latin
    2) run.font.name
    3) para.font.name
    4) para 的 pPr.rPr.rFonts/latin
    不再使用占位符、shape.lstStyle、母版textStyles、主题fontScheme等回退。
    其它属性（字号/颜色/粗斜体/下划线/删除线）维持原有 run/para 级别解析。
    """
    props: Dict[str, Any] = {}
    try:
        rf = run.font
        pf = getattr(para, 'font', None)
        # 字体名解析仅限：run.rPr / run.font / para.font / para.pPr
        name = None
        name_src = None
        reason = None
        try:
            rPr = run._r.rPr if hasattr(run, '_r') else None
            if rPr is not None and getattr(rPr, 'rFonts', None) is not None:
                rfonts = rPr.rFonts
                # 依次尝试 eastAsia、ascii、hAnsi、cs、latin.typeface
                name = (
                    getattr(rfonts, 'eastAsia', None)
                    or getattr(rfonts, 'ascii', None)
                    or getattr(rfonts, 'hAnsi', None)
                    or getattr(rfonts, 'cs', None)
                )
                if name is None and getattr(rPr, 'latin', None) is not None:
                    name = getattr(rPr.latin, 'typeface', None)
                if name is not None:
                    name_src = 'run.rPr.rFonts/latin'
        except Exception:
            pass
        if name is None and getattr(rf, 'name', None) is not None:
            name = rf.name
            name_src = 'run.font.name'
        if name is None and pf is not None and getattr(pf, 'name', None) is not None:
            name = pf.name
            name_src = 'para.font.name'
        if name is None:
            val = _get_para_rfonts(para)
            if val is not None:
                name = val
                name_src = 'para.pPr.rPr.rFonts/latin'
        if name is None:
            reason = 'run/paragraph 层均未提供字体名'
        # 规范化字体名称（去除前后空白）；如果所有方法都无法获取字体名，则设为"未知"
        if isinstance(name, str):
            name = name.strip()
            if name.startswith('+'):
                # 尝试从主题中获取占位符对应的字体
                try:
                    # 需要从shape获取presentation对象
                    if hasattr(host_shape, 'slide') and hasattr(host_shape.slide, 'presentation'):
                        presentation = host_shape.slide.presentation
                        theme_font = _get_theme_font_for_placeholder(name, presentation)
                        if theme_font:
                            name = theme_font
                            name_src = f'{name_src} -> 主题解析({name})' if name_src else f'主题解析({name})'
                            reason = None
                        else:
                            name = "未知"
                            name_src = name_src + ' -> 未知' if name_src else '未知'
                            reason = '主题占位符解析失败'
                    else:
                        name = "未知"
                        name_src = name_src + ' -> 未知' if name_src else '未知'
                        reason = '无法获取presentation对象'
                except Exception as e:
                    name = "未知"
                    name_src = name_src + ' -> 未知' if name_src else '未知'
                    reason = f'主题解析异常: {e}'
        elif name is None:
            # 所有方法都无法获取字体名，设为"未知"
            name = "未知"
            name_src = '未知'
            if reason is None:
                reason = '未从 run/para 获取到字体名'
        # 字体族合并（规范别名/派生名）
        merged = _merge_font_family_alias(name)
        
        # 将不在指定5种字体类型内的字体归为"其他"
        allowed_fonts = {"Meiryo UI", "宋体", "微软雅黑", "楷体", "Time New Roman"}
        if merged not in allowed_fonts and merged != "未知":
            props["字体类型"] = "其他"
        else:
            props["字体类型"] = merged

        # 可选：未知时简单提示（保留最小化日志）
        try:
            if merged == "未知":
                sid = str(getattr(host_shape, "shape_id", ""))
                snippet = ''
                try:
                    snippet = (run.text or '')[:30]
                except Exception:
                    snippet = ''
                print(f"[字体类型=未知] shape_id={sid} 源={name_src} 原因={reason or '无'} 文本片段='{snippet}'")
        except Exception:
            pass
        # 字号
        size = getattr(rf, 'size', None) or (getattr(pf, 'size', None) if pf is not None else None)
        props["字号"] = float(size.pt) if size is not None else 18.0
        # 颜色
        color = getattr(rf, 'color', None)
        hex_color = _rgb_to_hex(color) if color is not None else None
        if hex_color is None and pf is not None:
            pcolor = getattr(pf, 'color', None)
            hex_color = _rgb_to_hex(pcolor) if pcolor is not None else None
        if hex_color is None:
            hex_color = "#000000"
        props["字体颜色"] = hex_color
        # 样式
        def _bool_or_default(val, fallback):
            return bool(val) if val is not None else fallback
        bold = getattr(rf, 'bold', None)
        if bold is None and pf is not None:
            bold = getattr(pf, 'bold', None)
        italic = getattr(rf, 'italic', None)
        if italic is None and pf is not None:
            italic = getattr(pf, 'italic', None)
        underline = getattr(rf, 'underline', None)
        if underline is None and pf is not None:
            underline = getattr(pf, 'underline', None)
        strike = getattr(rf, 'strike', None)
        if strike is None and pf is not None:
            strike = getattr(pf, 'strike', None)
        props.update({
            "是否粗体": _bool_or_default(bold, False),
            "是否斜体": _bool_or_default(italic, False),
            "是否下划线": _bool_or_default(underline, False),
            "是否带删除线": _bool_or_default(strike, False),
        })
    except Exception:
        pass
    return props


def _process_table_cell(cell, cell_index: int, position: Dict[str, str]) -> Dict[str, Any]:
    """处理表格单元格的文本内容"""
    cell_text_info = {}
    try:
        if hasattr(cell, 'text_frame') and cell.text_frame:
            # 构建单元格文本块数据
            cell_text_data = {
                "文本块位置": position,
                "图层编号": cell_index,
                "是否是标题占位符": False,
                "文本块索引": f"table_cell_{cell_index}",
                "段落属性": []
            }
            
            # 获取单元格文本框架中的段落和运行
            text_frame = cell.text_frame
            for para_index, paragraph in enumerate(text_frame.paragraphs):
                for run_index, run in enumerate(paragraph.runs):
                    if run.text.strip():  # 只处理有文本的运行
                        # 构建运行属性对象
                        char_attr = {
                            "段落编号": para_index,
                            "字体类型": "未知",  # 表格单元格字体信息可能不完整
                            "字号": 12.0,  # 默认字号
                            "字体颜色": "黑色",  # 默认颜色
                            "是否粗体": False,
                            "段落内容": run.text
                        }
                        
                        # 使用与普通文本块相同的合并逻辑
                        if cell_text_data["段落属性"]:
                            last = cell_text_data["段落属性"][-1]
                            same_style = all(last.get(k) == char_attr.get(k) for k in ATTR_COMPARE_KEYS)
                            same_para = last.get("段落编号") == char_attr.get("段落编号")
                            if same_style and same_para:
                                last["段落内容"] = f"{last.get('段落内容','')}{char_attr.get('段落内容','')}"
                            else:
                                cell_text_data["段落属性"].append(char_attr)
                        else:
                            cell_text_data["段落属性"].append(char_attr)
            
            # 只有当有内容时才输出
            if cell_text_data["段落属性"]:
                text_key = f"文本块{cell_index + 1}"
                cell_text_info[text_key] = cell_text_data
                
    except Exception as e:
        print(f"处理表格单元格失败: {e}")
    
    return cell_text_info


def _get_text_block_info(shape, shape_index: int) -> Dict[str, Any]:
    text_info = {}
    try:
        # 检查是否为组合元素
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # 递归处理组合元素中的子形状
            group_text_info = {}
            for i, sub_shape in enumerate(shape.shapes):
                # 使用数字索引，避免字符串拼接问题
                sub_text_info = _get_text_block_info(sub_shape, shape_index * 100 + i)
                if sub_text_info:
                    group_text_info.update(sub_text_info)
            return group_text_info
        
        # 处理表格形状
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_text_info = {}
            try:
                table = shape.table
                # 计算幻灯片尺寸
                try:
                    slide_width = shape.part.slide_width
                    slide_height = shape.part.slide_height
                except Exception:
                    slide_width = 9144000
                    slide_height = 6858000
                
                # EMU 转 百分比
                def emu_to_percent_str(emu_val, slide_dimension) -> str:
                    try:
                        percent = (float(emu_val) / float(slide_dimension)) * 100.0
                        return f"{percent:.2f}%"
                    except Exception:
                        return "0.00%"
                
                # 遍历表格的每个单元格
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text.strip():  # 只处理有文本的单元格
                            # 计算单元格的绝对位置（基于表格左上角 + 行高/列宽累计）
                            try:
                                # 列偏移与宽度
                                col_offset = 0
                                for c in range(col_idx):
                                    col_offset += table.columns[c].width
                                col_width = table.columns[col_idx].width
                                # 行偏移与高度
                                row_offset = 0
                                for r in range(row_idx):
                                    row_offset += table.rows[r].height
                                row_height = table.rows[row_idx].height
                                # 形状左上角 + 偏移
                                cell_left = shape.left + col_offset
                                cell_top = shape.top + row_offset
                                # 百分比位置
                                position = {
                                    "left": emu_to_percent_str(cell_left, slide_width),
                                    "top": emu_to_percent_str(cell_top, slide_height),
                                    "width": emu_to_percent_str(col_width, slide_width),
                                    "height": emu_to_percent_str(row_height, slide_height),
                                }
                            except Exception:
                                # 回退为整表位置
                                pos = _get_shape_position(shape)
                                position = {"left": pos.get("left", "0%"), "top": pos.get("top", "0%"),
                                            "width": pos.get("width", "100%"), "height": pos.get("height", "100%")}

                            # 表格单元格有text_frame，直接处理文本内容
                            cell_text_info = _process_table_cell(
                                cell,
                                shape_index * 1000 + row_idx * 100 + col_idx,
                                position,
                            )
                            if cell_text_info:
                                table_text_info.update(cell_text_info)
                
                return table_text_info
            except Exception as e:
                print(f"处理表格形状失败: {e}")
                return {}
        
        # 处理普通文本形状
        if shape.has_text_frame and shape.text_frame:
            text_block_position = _get_shape_position(shape)
            is_title_placeholder = _is_title_placeholder(shape)
            
            # 构建文本块数据（改：输出“段落属性”，不再生成“拼接字符”）
            text_block_data = {
                "文本块位置": text_block_position,
                "图层编号": shape_index,
                "是否是标题占位符": is_title_placeholder,
                "文本块索引": str(getattr(shape, "shape_id", "")),
                "段落属性": []
            }
            
            # 获取文本框架中的段落和运行
            text_frame = shape.text_frame
            for para_index, paragraph in enumerate(text_frame.paragraphs):
                # 处理段落中的每个运行（run）
                for run_index, run in enumerate(paragraph.runs):
                    # 获取运行的字体属性
                    font_props = _resolve_run_font_props(run, paragraph, is_title_placeholder, shape)
                    
                    # 构建 run 级属性对象（段落属性项）
                    char_attr = {
                        "段落编号": para_index,
                        "字体类型": font_props.get("字体类型", "未知"),
                        "字号": font_props.get("字号", 18.0),
                        # 颜色改为中文常见色名
                        "字体颜色": _hex_to_cn_color_name(font_props.get("字体颜色", "#000000")),
                        "是否粗体": font_props.get("是否粗体", False),
                        # "是否斜体": font_props.get("是否斜体", False),
                        # "是否下划线": font_props.get("是否下划线", False),
                        # "是否带删除线": font_props.get("是否带删除线", False),
                        "段落内容": run.text
                    }

                    # JSON层面合并：若与前一条在样式上完全一致（仅段落内容不同），则合并段落内容
                    if text_block_data["段落属性"]:
                        last = text_block_data["段落属性"][-1]
                        same_style = all(last.get(k) == char_attr.get(k) for k in ATTR_COMPARE_KEYS)
                        same_para = last.get("段落编号") == char_attr.get("段落编号")
                        if same_style and same_para:
                            last["段落内容"] = f"{last.get('段落内容','')}{char_attr.get('段落内容','')}"
                        else:
                            text_block_data["段落属性"].append(char_attr)
                    else:
                        text_block_data["段落属性"].append(char_attr)

            # 检查是否有有效的段落内容
            has_content = False
            for char_attr in text_block_data["段落属性"]:
                if char_attr.get("段落内容", "").strip():
                    has_content = True
                    break
            
            # 只有当有内容时才输出文本块
            if has_content:
                text_key = f"文本块{shape_index + 1}"
                sid = str(getattr(shape, "shape_id", ""))
                text_payload = {
                    "文本块位置": text_block_position,
                    "图层编号": shape_index,
                    "是否是标题占位符": is_title_placeholder,
                    "文本块索引": sid,
                    "段落属性": text_block_data["段落属性"]
                }
                text_info = {text_key: text_payload}
    except Exception as e:
        print(f"提取文本块信息失败: {e}")
    return text_info


def _get_image_info(shape, shape_index: int) -> Dict[str, Any]:
    image_info = {}
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            position = _get_shape_position(shape)
            image_type = "未知"
            try:
                if hasattr(shape, 'image') and shape.image:
                    filename = getattr(shape.image, 'filename', '')
                    if filename:
                        image_type = filename.split('.')[-1].upper() if '.' in filename else "未知"
            except Exception:
                pass
            image_size = 0
            try:
                if hasattr(shape, 'image') and shape.image:
                    image_size = len(shape.image.blob) if hasattr(shape.image, 'blob') else 0
            except Exception:
                pass
            # 提取 shape_id 作为图片索引（便于与原始PPT形状对应）
            sid = str(getattr(shape, "shape_id", ""))
            image_info = {
                f"图片{shape_index + 1}": {
                    "图片位置": position,
                    "图片类型": image_type,
                    "图片大小": f"{image_size} bytes",
                    "图层位置": shape_index,
                    "图片索引": sid
                }
            }
    except Exception as e:
        print(f"提取图片信息失败: {e}")
    return image_info


def parse_pptx(path: str, include_images: bool = False) -> Dict[str, Any]:
    try:
        prs = Presentation(path)
        total_slides = len(prs.slides)
        
        # 按照新结构组织数据
        result = {
            "页数": total_slides,
            "contents": []
        }
        
        for slide_index, slide in enumerate(prs.slides):
            page_data = {
                "页码": slide_index + 1,
                "文本块数量": 0,
                "文本块": [],
                "图片数量": 0,
                "图片": []
            }
            
            # 处理文本块（包括组合元素中的文本）
            text_blocks: List[Dict[str, Any]] = []
            for shape_index, shape in enumerate(slide.shapes):
                text_info = _get_text_block_info(shape, shape_index)
                if text_info:
                    # 提取文本块内容到数组
                    for key, payload in text_info.items():
                        if key.startswith("文本块"):
                            text_blocks.append(payload)
            
            page_data["文本块数量"] = len(text_blocks)
            page_data["文本块"] = text_blocks
            
            # 处理图片
            if include_images:
                images: List[Dict[str, Any]] = []
                for shape_index, shape in enumerate(slide.shapes):
                    image_info = _get_image_info(shape, shape_index)
                    if image_info:
                        # 提取图片内容到数组
                        for key, payload in image_info.items():
                            if key.startswith("图片"):
                                images.append(payload)
                
                page_data["图片数量"] = len(images)
                page_data["图片"] = images
            
            result["contents"].append(page_data)
        
        return result
    except Exception as e:
        print(f"解析PPTX文件失败: {e}")
        return {"页数": 0, "contents": []}


def save_to_json(data: List[Dict[str, Any]], output_path: str):
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        print(f"数据已保存到: {output_path}")
    except Exception as e:
        print(f"保存JSON文件失败: {e}")


if __name__ == "__main__":
    # 命令行参数：--include-images 控制是否输出图片信息（默认不输出）
    import argparse
    parser = argparse.ArgumentParser(description="PPTX解析器")
    parser.add_argument("pptx", nargs='?', default="example2.pptx", help="PPTX 文件路径")
    parser.add_argument("--include-images", action="store_true", help="是否输出图片信息，默认否")
    args = parser.parse_args()

    pptx_path = args.pptx
    print("🧪 测试PPTX解析...")
    # 依据新增参数 include_images 控制图片信息输出
    result = parse_pptx(pptx_path, include_images=args.include_images)
    if result and "contents" in result:
        print(f"✅ 成功解析，共 {result['页数']} 页")
        save_to_json(result, "parsing_result.json")
        
        # # 显示前几页的关键信息
        # for i, page in enumerate(result['contents'][:5]):
        #     print(f"\n第 {page['页码']} 页:")
            
        #     # 显示前几个文本块的关键信息
        #     for j, text_block in enumerate(page.get('文本块', [])):  # 只显示前3个文本块
        #         print(text_block)
    else:
        print("❌ 解析失败")

