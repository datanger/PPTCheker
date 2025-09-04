"""
å·¥ä½œæµå·¥å…·æ¨¡å—ï¼šæä¾›å·¥ä½œæµæ‰€éœ€çš„æ‰€æœ‰å·¥å…·å‡½æ•°

åŠŸèƒ½ï¼š
1. ä»Ž parsing_result.json åŠ è½½å’Œè§£æžæ•°æ®
2. æ‰§è¡ŒåŸºç¡€è§„åˆ™æ£€æŸ¥
3. æ‰§è¡ŒLLMæ™ºèƒ½å®¡æŸ¥
4. ç”Ÿæˆå®¡æŸ¥æŠ¥å‘Š
5. ç”Ÿæˆæ ‡è®°PPT

è¾“å…¥ï¼šparsing_result.json æ ¼å¼çš„æ•°æ®
è¾“å‡ºï¼šå®¡æŸ¥ç»“æžœã€æŠ¥å‘Šã€æ ‡è®°PPTç­‰
"""

import json
import os
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path

# å¯¼å…¥å¿…è¦çš„æ¨¡å—
try:
    from ..model import Issue, DocumentModel, Slide, Shape, TextRun, PPTContext, EditSuggestion, EditResult
    from ..config import ToolConfig
    from ..llm import LLMClient
    from ..reporter import render_markdown
    from ..annotator import annotate_pptx
except ImportError:
    # å…¼å®¹ç›´æŽ¥è¿è¡Œçš„æƒ…å†µ
    import sys
    sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    from model import Issue, DocumentModel, Slide, Shape, TextRun
    from config import ToolConfig
    from llm import LLMClient
    from reporter import render_markdown
    from annotator import annotate_pptx


def load_parsing_result(file_path: str = "parsing_result.json") -> Dict[str, Any]:
    """åŠ è½½ parsing_result.json æ–‡ä»¶"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        print(f"âŒ åŠ è½½ parsing_result.json å¤±è´¥: {e}")
        return {"é¡µæ•°": 0, "contents": []}


def convert_parsing_result_to_document_model(parsing_data: Dict[str, Any], file_path: str) -> DocumentModel:
    """å°† parsing_result.json æ ¼å¼è½¬æ¢ä¸º DocumentModel æ ¼å¼"""
    slides = []
    
    for page_data in parsing_data.get("contents", []):
        slide = Slide(
            index=page_data["é¡µç "] - 1,  # è½¬æ¢ä¸º0åŸºç´¢å¼•
            slide_title=page_data.get("é¡µæ ‡é¢˜", ""),
            slide_type=page_data.get("é¡µç±»åž‹", "å†…å®¹é¡µ"),
            chapter_info=None
        )
        
        # å¤„ç†æ–‡æœ¬å—ï¼ˆæ”¹ï¼šä½¿ç”¨â€œæ®µè½å±žæ€§â€ç›´æŽ¥æž„å»º TextRunï¼Œå·²ä¸å†ä¾èµ–â€œæ‹¼æŽ¥å­—ç¬¦â€ï¼‰
        for text_block in page_data.get("æ–‡æœ¬å—", []):
            shape = Shape(
                id=str(text_block["æ–‡æœ¬å—ç´¢å¼•"]),
                type="text",
                is_title=text_block.get("æ˜¯å¦æ˜¯æ ‡é¢˜å ä½ç¬¦", False),
                title_level=1 if text_block.get("æ˜¯å¦æ˜¯æ ‡é¢˜å ä½ç¬¦", False) else None,
                text_color=None,
                fill_color=None,
                border_color=None
            )
            
            # ç›´æŽ¥ä»Ž"æ®µè½å±žæ€§"æž„å»º TextRun
            para_runs = text_block.get("æ®µè½å±žæ€§", [])
            for r in para_runs:
                tr = TextRun(
                    text=str(r.get("æ®µè½å†…å®¹", "")),
                    font_name=r.get("å­—ä½“ç±»åž‹"),
                    font_size_pt=float(r.get("å­—å·")) if r.get("å­—å·") is not None else None,
                    language_tag="ja",
                    is_bold=bool(r.get("æ˜¯å¦ç²—ä½“", False)),
                    is_italic=bool(r.get("æ˜¯å¦æ–œä½“", False)),
                    is_underline=bool(r.get("æ˜¯å¦ä¸‹åˆ’çº¿", False))
                )
                shape.text_runs.append(tr)
            
            slide.shapes.append(shape)
        
        slides.append(slide)
    
    doc = DocumentModel(file_path=file_path, slides=slides)
    return doc


def _parse_concatenated_text_to_runs(concatenated_text: str) -> List[TextRun]:
    """ä»Žæ‹¼æŽ¥å­—ç¬¦ä¸­è§£æžå‡ºæ–‡æœ¬è¿è¡Œ"""
    runs = []
    
    # ç®€å•çš„æ–‡æœ¬åˆ†å‰²ï¼ˆå®žé™…åº”è¯¥æ›´æ™ºèƒ½åœ°è§£æžï¼‰
    # è¿™é‡Œç®€åŒ–å¤„ç†ï¼Œå°†æ•´ä¸ªæ–‡æœ¬ä½œä¸ºä¸€ä¸ªrun
    if concatenated_text:
        # æå–çº¯æ–‡æœ¬å†…å®¹ï¼ˆåŽ»é™¤å±žæ€§æ ‡è®°ï¼‰
        clean_text = _extract_clean_text(concatenated_text)
        
        run = TextRun(
            text=clean_text,
            font_name="é»˜è®¤å­—ä½“",  # ä»Žå±žæ€§ä¸­æå–
            font_size_pt=18.0,     # ä»Žå±žæ€§ä¸­æå–
            language_tag="ja",      # é»˜è®¤æ—¥è¯­
            is_bold=False,
            is_italic=False,
            is_underline=False
        )
        runs.append(run)
    
    return runs


def _extract_clean_text(concatenated_text: str) -> str:
    """ä»Žæ‹¼æŽ¥å­—ç¬¦ä¸­æå–çº¯æ–‡æœ¬å†…å®¹"""
    # ç§»é™¤æ‰€æœ‰å±žæ€§æ ‡è®°
    import re
    
    # ç§»é™¤ã€åˆå§‹çš„å­—ç¬¦æ‰€æœ‰å±žæ€§ï¼š...ã€‘æ ‡è®°
    text = re.sub(r'ã€åˆå§‹çš„å­—ç¬¦æ‰€æœ‰å±žæ€§ï¼š[^ã€‘]*ã€‘', '', concatenated_text)
    
    # ç§»é™¤ã€å­—ç¬¦å±žæ€§å˜æ›´ï¼š...ã€‘æ ‡è®°
    text = re.sub(r'ã€å­—ç¬¦å±žæ€§å˜æ›´ï¼š[^ã€‘]*ã€‘', '', text)
    
    # ç§»é™¤ã€æ¢è¡Œã€‘æ ‡è®°
    text = re.sub(r'ã€æ¢è¡Œã€‘', '\n', text)
    
    # ç§»é™¤ã€ç¼©è¿›{...}ã€‘æ ‡è®°
    text = re.sub(r'ã€ç¼©è¿›\{\d+\}ã€‘', '', text)
    
    return text.strip()


def run_basic_rules(doc: DocumentModel, cfg: ToolConfig) -> List[Issue]:
    """è¿è¡ŒåŸºç¡€è§„åˆ™æ£€æŸ¥"""
    from .rules import run_basic_rules as run_rules
    return run_rules(doc, cfg)


def run_llm_review(doc: DocumentModel, llm: LLMClient, cfg: ToolConfig) -> List[Issue]:
    """è¿è¡ŒLLMæ™ºèƒ½å®¡æŸ¥"""
    try:
        from .llm_review import create_llm_reviewer
        reviewer = create_llm_reviewer(llm, cfg)
        return reviewer.run_llm_review(doc)
    except Exception as e:
        print(f"âš ï¸ LLMå®¡æŸ¥å¤±è´¥: {e}")
        return []


def generate_report(issues: List[Issue], rule_issues: List[Issue] = None, llm_issues: List[Issue] = None) -> str:
    """ç”Ÿæˆå®¡æŸ¥æŠ¥å‘Š"""
    try:
        # å¦‚æžœæ²¡æœ‰æä¾›åˆ†ç±»ä¿¡æ¯ï¼Œä½¿ç”¨é»˜è®¤çš„render_markdown
        if rule_issues is None or llm_issues is None:
            return render_markdown(issues)
        
        # ç”Ÿæˆåˆ†ç±»æŠ¥å‘Š
        return _generate_categorized_report(issues, rule_issues, llm_issues)
    except Exception as e:
        print(f"âš ï¸ ç”ŸæˆæŠ¥å‘Šå¤±è´¥: {e}")
        return f"# å®¡æŸ¥æŠ¥å‘Š\n\nç”ŸæˆæŠ¥å‘Šæ—¶å‘ç”Ÿé”™è¯¯: {e}"


def _generate_categorized_report(issues: List[Issue], rule_issues: List[Issue], llm_issues: List[Issue]) -> str:
    """ç”Ÿæˆåˆ†ç±»æŠ¥å‘Š"""
    # åˆ›å»ºè§„åˆ™æ£€æŸ¥å’ŒLLMå®¡æŸ¥çš„é—®é¢˜é›†åˆ
    rule_issue_ids = {id(issue) for issue in rule_issues}
    llm_issue_ids = {id(issue) for issue in llm_issues}
    
    # åˆ†ç±»é—®é¢˜
    categorized_rule_issues = []
    categorized_llm_issues = []
    
    for issue in issues:
        if id(issue) in rule_issue_ids:
            categorized_rule_issues.append(issue)
        elif id(issue) in llm_issues:
            categorized_llm_issues.append(issue)
        else:
            # å¦‚æžœæ— æ³•ç¡®å®šæ¥æºï¼Œæ ¹æ®rule_idåˆ¤æ–­
            if issue.rule_id.startswith("LLM_") or issue.rule_id.endswith("_AcronymRule") or issue.rule_id.endswith("_ContentRule") or issue.rule_id.endswith("_FormatRule") or issue.rule_id.endswith("_TitleStructureRule"):
                categorized_llm_issues.append(issue)
            else:
                categorized_rule_issues.append(issue)
    
    # ç”ŸæˆæŠ¥å‘Šå†…å®¹
    report = "# å®¡æŸ¥æŠ¥å‘Š\n\n"
    
    # é—®é¢˜ç»Ÿè®¡
    report += f"### ðŸ“Š é—®é¢˜ç»Ÿè®¡\n"
    report += f"- **è§„åˆ™æ£€æŸ¥é—®é¢˜**: {len(categorized_rule_issues)} ä¸ª\n"
    report += f"- **LLMæ™ºèƒ½å®¡æŸ¥é—®é¢˜**: {len(categorized_llm_issues)} ä¸ª\n"
    report += f"- **æ€»è®¡**: {len(issues)} ä¸ª\n\n"
    
    # æŒ‰é¡µç åˆ†ç»„æ˜¾ç¤ºé—®é¢˜
    report += "### ðŸ“„ æŒ‰é¡µç åˆ†ç»„çš„é—®é¢˜è¯¦æƒ…\n\n"
    
    # èŽ·å–æ‰€æœ‰é¡µç å¹¶æŽ’åº
    all_page_numbers = set()
    for issue in issues:
        page_num = issue.slide_index + 1  # è½¬æ¢ä¸º1åŸºé¡µç 
        all_page_numbers.add(page_num)
    
    if all_page_numbers:
        # æŒ‰é¡µç æŽ’åº
        sorted_pages = sorted(all_page_numbers)
        
        for page_num in sorted_pages:
            report += f"#### ðŸ“ ç¬¬ {page_num} é¡µ\n\n"
            
            # æ”¶é›†è¯¥é¡µçš„æ‰€æœ‰é—®é¢˜
            page_issues = [issue for issue in issues if issue.slide_index + 1 == page_num]
            
            if page_issues:
                # æŒ‰é—®é¢˜ç±»åž‹åˆ†ç»„
                rule_issues_on_page = [issue for issue in page_issues if not issue.rule_id.startswith("LLM_")]
                llm_issues_on_page = [issue for issue in page_issues if issue.rule_id.startswith("LLM_")]
                
                # æ˜¾ç¤ºè§„åˆ™æ£€æŸ¥é—®é¢˜
                if rule_issues_on_page:
                    report += "**ðŸ” è§„åˆ™æ£€æŸ¥é—®é¢˜:**\n\n"
                    for issue in rule_issues_on_page:
                        report += f"- **{issue.rule_id}** | ä¸¥é‡æ€§: {issue.severity} | å¯¹è±¡: {issue.object_ref}\n"
                        report += f"  - æè¿°: {issue.message}\n"
                        if issue.suggestion:
                            report += f"  - å»ºè®®: {issue.suggestion}\n"
                        report += f"  - å¯è‡ªåŠ¨ä¿®å¤: {'æ˜¯' if issue.can_autofix else 'å¦'} | å·²ä¿®å¤: {'æ˜¯' if getattr(issue, 'is_fixed', False) else 'å¦'}\n\n"
                
                # æ˜¾ç¤ºLLMå®¡æŸ¥é—®é¢˜
                if llm_issues_on_page:
                    report += "**ðŸ¤– LLMæ™ºèƒ½å®¡æŸ¥é—®é¢˜:**\n\n"
                    for issue in llm_issues_on_page:
                        report += f"- **{issue.rule_id}** | ä¸¥é‡æ€§: {issue.severity} | å¯¹è±¡: {issue.object_ref}\n"
                        report += f"  - æè¿°: {issue.message}\n"
                        if issue.suggestion:
                            report += f"  - å»ºè®®: {issue.suggestion}\n"
                        report += f"  - å¯è‡ªåŠ¨ä¿®å¤: {'æ˜¯' if issue.can_autofix else 'å¦'} | å·²ä¿®å¤: {'æ˜¯' if getattr(issue, 'is_fixed', False) else 'å¦'}\n\n"
                
                # æ˜¾ç¤ºè¯¥é¡µé—®é¢˜ç»Ÿè®¡
                report += f"**ðŸ“Š ç¬¬ {page_num} é¡µé—®é¢˜ç»Ÿè®¡:** å…± {len(page_issues)} ä¸ªé—®é¢˜\n\n"
            else:
                report += "âœ… è¯¥é¡µæœªå‘çŽ°é—®é¢˜\n\n"
            
            report += "---\n\n"
    else:
        report += "âœ… æœªå‘çŽ°ä»»ä½•é—®é¢˜\n\n"
    
    # é—®é¢˜åˆ†ç±»ç»Ÿè®¡
    report += "### ðŸ“‹ é—®é¢˜åˆ†ç±»ç»Ÿè®¡\n"
    
    # è§„åˆ™æ£€æŸ¥åˆ†ç±»
    if categorized_rule_issues:
        rule_counts = {}
        for issue in categorized_rule_issues:
            rule_counts[issue.rule_id] = rule_counts.get(issue.rule_id, 0) + 1
        
        report += "**è§„åˆ™æ£€æŸ¥åˆ†ç±»:**\n\n"
        for rule_id, count in rule_counts.items():
            report += f"- {rule_id}: {count} ä¸ª\n"
    else:
        report += "**è§„åˆ™æ£€æŸ¥åˆ†ç±»:**\n\næ— \n"
    
    # LLMå®¡æŸ¥åˆ†ç±»
    if categorized_llm_issues:
        llm_counts = {}
        for issue in categorized_llm_issues:
            llm_counts[issue.rule_id] = llm_counts.get(issue.rule_id, 0) + 1
        
        report += "\n**LLMå®¡æŸ¥åˆ†ç±»:**\n\n"
        for rule_id, count in llm_counts.items():
            report += f"- {rule_id}: {count} ä¸ª\n"
    else:
        report += "\n**LLMå®¡æŸ¥åˆ†ç±»:**\n\næ— \n"
    
    return report


def generate_annotated_ppt(input_ppt: str, issues: List[Issue], output_ppt: str) -> bool:
    """ç”Ÿæˆæ ‡è®°PPT"""
    try:
        annotate_pptx(input_ppt, issues, output_ppt)
        return True
    except Exception as e:
        print(f"âš ï¸ ç”Ÿæˆæ ‡è®°PPTå¤±è´¥: {e}")
        return False


def get_workflow_statistics(rule_issues: List[Issue], llm_issues: List[Issue]) -> Dict[str, Any]:
    """èŽ·å–å·¥ä½œæµç»Ÿè®¡ä¿¡æ¯"""
    return {
        "rule_issues_count": len(rule_issues),
        "llm_issues_count": len(llm_issues),
        "total_issues": len(rule_issues) + len(llm_issues),
        "issues_by_severity": _count_issues_by_severity(rule_issues + llm_issues),
        "issues_by_rule": _count_issues_by_rule(rule_issues + llm_issues)
    }


# æ–°å¢žï¼šPPTç¼–è¾‘ç›¸å…³åŠŸèƒ½
def create_ppt_context(parsing_data: Dict[str, Any], original_pptx_path: str) -> Optional[PPTContext]:
    """åˆ›å»ºPPTç¼–è¾‘ä¸Šä¸‹æ–‡"""
    try:
        from pptx import Presentation
        
        # åŠ è½½åŽŸå§‹PPT
        prs = Presentation(original_pptx_path)
        
        # æå–ä¸»é¢˜ä¿¡æ¯
        theme_info = extract_theme_info(prs)
        
        # åˆ›å»ºä¸Šä¸‹æ–‡å¯¹è±¡
        context = PPTContext(
            parsing_result=parsing_data,
            original_pptx_path=original_pptx_path,
            presentation_object=prs,
            slide_layouts=list(prs.slide_layouts),
            slide_masters=list(prs.slide_masters),
            theme_info=theme_info
        )
        
        print(f"âœ… æˆåŠŸåˆ›å»ºPPTç¼–è¾‘ä¸Šä¸‹æ–‡ï¼Œå…± {len(prs.slides)} é¡µ")
        return context
        
    except Exception as e:
        print(f"âš ï¸ åˆ›å»ºPPTä¸Šä¸‹æ–‡å¤±è´¥: {e}")
        return None


def extract_theme_info(prs) -> Dict[str, Any]:
    """æå–PPTä¸»é¢˜ä¿¡æ¯"""
    theme_info = {}
    try:
        # æå–ä¸»é¢˜è‰²
        if hasattr(prs, 'core_properties'):
            theme_info['title'] = getattr(prs.core_properties, 'title', '')
            theme_info['author'] = getattr(prs.core_properties, 'author', '')
            theme_info['created'] = getattr(prs.core_properties, 'created', '')
        
        # æå–æ¯ç‰ˆä¿¡æ¯
        if hasattr(prs, 'slide_masters'):
            theme_info['slide_masters_count'] = len(prs.slide_masters)
        
        # æå–å¸ƒå±€ä¿¡æ¯
        if hasattr(prs, 'slide_layouts'):
            theme_info['slide_layouts_count'] = len(prs.slide_layouts)
            
    except Exception as e:
        print(f"âš ï¸ æå–ä¸»é¢˜ä¿¡æ¯å¤±è´¥: {e}")
    
    return theme_info


def run_llm_edit_analysis(parsing_data: Dict[str, Any], llm: LLMClient, edit_requirements: str) -> List[EditSuggestion]:
    """ä½¿ç”¨LLMåˆ†æžå¹¶ç”Ÿæˆç¼–è¾‘å»ºè®®"""
    try:
        # æž„å»ºç¼–è¾‘åˆ†æžæç¤ºè¯
        prompt = f"""
            ä½ æ˜¯PPTç¼–è¾‘ä¸“å®¶ã€‚åŸºäºŽä»¥ä¸‹PPTå†…å®¹åˆ†æžï¼Œè¯·æä¾›å…·ä½“çš„ç¼–è¾‘å»ºè®®ï¼š

            PPTå†…å®¹ï¼š
            {json.dumps(parsing_data, ensure_ascii=False, indent=2)}

            ç¼–è¾‘è¦æ±‚ï¼š
            {edit_requirements}

            è¯·åˆ†æžPPTå†…å®¹ï¼Œè¯†åˆ«éœ€è¦æ”¹è¿›çš„åœ°æ–¹ï¼Œå¹¶è¾“å‡ºJSONæ ¼å¼çš„ç¼–è¾‘å»ºè®®ã€‚

            è¾“å‡ºæ ¼å¼ï¼ˆåªè¾“å‡ºJSONæ•°ç»„ï¼Œä¸è¦è§£é‡Šï¼‰ï¼š
            [
            {{
                "type": "text_change|font_change|color_change|layout_change",
                "page_number": 1,
                "shape_index": 0,
                "current_value": "å½“å‰å€¼",
                "new_value": "æ–°å€¼",
                "reason": "ä¿®æ”¹åŽŸå› ",
                "priority": "high|medium|low",
                "can_auto_apply": true
            }}
            ]

            æ³¨æ„ï¼š
            1. page_number ä»Ž1å¼€å§‹è®¡æ•°
            2. shape_index æ˜¯è¯¥é¡µä¸­å½¢çŠ¶çš„ç´¢å¼•ï¼ˆä»Ž0å¼€å§‹ï¼‰
            3. åªæä¾›ç¡®å®žéœ€è¦ä¿®æ”¹çš„å»ºè®®
            4. ç¡®ä¿å»ºè®®å…·ä½“ä¸”å¯æ‰§è¡Œ
            """
                    
        # è°ƒç”¨LLM
        response = llm.complete(prompt=prompt, max_tokens=2048)
        
        # è§£æžJSONå“åº”
        try:
            suggestions_data = json.loads(response.strip())
            suggestions = []
            
            for item in suggestions_data:
                suggestion = EditSuggestion(
                    type=item.get('type', 'text_change'),
                    page_number=item.get('page_number', 1),
                    shape_index=item.get('shape_index', 0),
                    current_value=item.get('current_value', ''),
                    new_value=item.get('new_value', ''),
                    reason=item.get('reason', ''),
                    priority=item.get('priority', 'medium'),
                    can_auto_apply=item.get('can_auto_apply', True)
                )
                suggestions.append(suggestion)
            
            print(f"âœ… LLMç”Ÿæˆ {len(suggestions)} ä¸ªç¼–è¾‘å»ºè®®")
            return suggestions
            
        except json.JSONDecodeError as e:
            print(f"âš ï¸ è§£æžLLMå“åº”å¤±è´¥: {e}")
            print(f"LLMåŽŸå§‹å“åº”: {response}")
            return []
            
    except Exception as e:
        print(f"âš ï¸ LLMç¼–è¾‘åˆ†æžå¤±è´¥: {e}")
        return []


def apply_edits_to_ppt(ppt_context: PPTContext, edit_suggestions: List[EditSuggestion]) -> EditResult:
    """åº”ç”¨ç¼–è¾‘å»ºè®®åˆ°PPT"""
    result = EditResult(success=False)
    
    if not ppt_context or not ppt_context.presentation_object:
        result.error_messages.append("PPTä¸Šä¸‹æ–‡æ— æ•ˆ")
        return result
    
    try:
        for suggestion in edit_suggestions:
            try:
                # èŽ·å–ç›®æ ‡å¹»ç¯ç‰‡
                slide = ppt_context.get_editable_slide(suggestion.page_number)
                if not slide:
                    result.failed_suggestions.append(suggestion)
                    result.error_messages.append(f"é¡µé¢ {suggestion.page_number} ä¸å­˜åœ¨")
                    continue
                
                # èŽ·å–ç›®æ ‡å½¢çŠ¶
                if suggestion.shape_index >= len(slide.shapes):
                    result.failed_suggestions.append(suggestion)
                    result.error_messages.append(f"é¡µé¢ {suggestion.page_number} çš„å½¢çŠ¶ {suggestion.shape_index} ä¸å­˜åœ¨")
                    continue
                
                shape = slide.shapes[suggestion.shape_index]
                
                # æ ¹æ®ç±»åž‹åº”ç”¨ç¼–è¾‘
                if suggestion.type == "text_change":
                    # ä¿®æ”¹æ–‡æœ¬
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        shape.text_frame.text = suggestion.new_value
                        result.applied_suggestions.append(suggestion)
                        result.modified_slides.append(suggestion.page_number)
                        print(f"âœ… é¡µé¢ {suggestion.page_number} å½¢çŠ¶ {suggestion.shape_index} æ–‡æœ¬å·²ä¿®æ”¹")
                    else:
                        result.failed_suggestions.append(suggestion)
                        result.error_messages.append(f"å½¢çŠ¶ {suggestion.shape_index} ä¸æ”¯æŒæ–‡æœ¬ç¼–è¾‘")
                
                elif suggestion.type == "font_change":
                    # ä¿®æ”¹å­—ä½“
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = suggestion.new_value
                        result.applied_suggestions.append(suggestion)
                        result.modified_slides.append(suggestion.page_number)
                        print(f"âœ… é¡µé¢ {suggestion.page_number} å½¢çŠ¶ {suggestion.shape_index} å­—ä½“å·²ä¿®æ”¹")
                    else:
                        result.failed_suggestions.append(suggestion)
                        result.error_messages.append(f"å½¢çŠ¶ {suggestion.shape_index} ä¸æ”¯æŒå­—ä½“ç¼–è¾‘")
                
                elif suggestion.type == "color_change":
                    # ä¿®æ”¹é¢œè‰²
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        from ..model import Color
                        # è§£æžé¢œè‰²å€¼ï¼ˆå‡è®¾æ ¼å¼ä¸º #RRGGBBï¼‰
                        if suggestion.new_value.startswith('#'):
                            r = int(suggestion.new_value[1:3], 16)
                            g = int(suggestion.new_value[3:5], 16)
                            b = int(suggestion.new_value[5:7], 16)
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = (r, g, b)
                            result.applied_suggestions.append(suggestion)
                            result.modified_slides.append(suggestion.page_number)
                            print(f"âœ… é¡µé¢ {suggestion.page_number} å½¢çŠ¶ {suggestion.shape_index} é¢œè‰²å·²ä¿®æ”¹")
                        else:
                            result.failed_suggestions.append(suggestion)
                            result.error_messages.append(f"æ— æ•ˆçš„é¢œè‰²æ ¼å¼: {suggestion.new_value}")
                    else:
                        result.failed_suggestions.append(suggestion)
                        result.error_messages.append(f"å½¢çŠ¶ {suggestion.shape_index} ä¸æ”¯æŒé¢œè‰²ç¼–è¾‘")
                
                else:
                    result.failed_suggestions.append(suggestion)
                    result.error_messages.append(f"ä¸æ”¯æŒçš„ç¼–è¾‘ç±»åž‹: {suggestion.type}")
                
            except Exception as e:
                result.failed_suggestions.append(suggestion)
                result.error_messages.append(f"åº”ç”¨ç¼–è¾‘å»ºè®®å¤±è´¥: {e}")
                print(f"âš ï¸ åº”ç”¨ç¼–è¾‘å»ºè®®å¤±è´¥: {e}")
        
        # åŽ»é‡ä¿®æ”¹çš„é¡µé¢
        result.modified_slides = list(set(result.modified_slides))
        result.success = len(result.applied_suggestions) > 0
        
        print(f"âœ… ç¼–è¾‘å®Œæˆï¼šæˆåŠŸ {len(result.applied_suggestions)} ä¸ªï¼Œå¤±è´¥ {len(result.failed_suggestions)} ä¸ª")
        return result
        
    except Exception as e:
        result.error_messages.append(f"åº”ç”¨ç¼–è¾‘è¿‡ç¨‹ä¸­å‘ç”Ÿé”™è¯¯: {e}")
        print(f"âš ï¸ åº”ç”¨ç¼–è¾‘è¿‡ç¨‹ä¸­å‘ç”Ÿé”™è¯¯: {e}")
        return result


def save_modified_ppt(ppt_context: PPTContext, output_path: str) -> bool:
    """ä¿å­˜ä¿®æ”¹åŽçš„PPT"""
    try:
        if ppt_context and ppt_context.presentation_object:
            ppt_context.presentation_object.save(output_path)
            print(f"âœ… ä¿®æ”¹åŽçš„PPTå·²ä¿å­˜åˆ°: {output_path}")
            return True
        else:
            print("âŒ PPTä¸Šä¸‹æ–‡æ— æ•ˆï¼Œæ— æ³•ä¿å­˜")
            return False
    except Exception as e:
        print(f"âš ï¸ ä¿å­˜PPTå¤±è´¥: {e}")
        return False


def _count_issues_by_severity(issues: List[Issue]) -> Dict[str, int]:
    """æŒ‰ä¸¥é‡ç¨‹åº¦ç»Ÿè®¡é—®é¢˜"""
    counts = {"error": 0, "warning": 0, "info": 0}
    for issue in issues:
        severity = getattr(issue, 'severity', 'warning')
        counts[severity] = counts.get(severity, 0) + 1
    return counts


def _count_issues_by_rule(issues: List[Issue]) -> Dict[str, int]:
    """æŒ‰è§„åˆ™ç±»åž‹ç»Ÿè®¡é—®é¢˜"""
    counts = {}
    for issue in issues:
        rule_id = getattr(issue, 'rule_id', 'unknown')
        counts[rule_id] = counts.get(rule_id, 0) + 1
    return counts
