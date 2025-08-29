#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
检查output.pptx文件中的标记
"""

from pptx import Presentation

def check_markings():
    """检查PPT文件中的标记"""
    try:
        # 加载标记后的PPT
        prs = Presentation('output.pptx')
        print(f"总页数: {len(prs.slides)}")
        
        # 搜索整个PPT中包含ADAS的形状
        adas_shapes = []
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame is not None:
                    all_text = ""
                    for para in shape.text_frame.paragraphs:
                        all_text += para.text
                    
                    if "ADAS" in all_text:
                        adas_shapes.append((slide_idx, shape_idx, shape, all_text))
        
        print(f"整个PPT包含ADAS的形状数量: {len(adas_shapes)}")
        
        # 显示包含ADAS的形状
        for slide_idx, shape_idx, shape, text in adas_shapes:
            print(f"\n页面{slide_idx + 1} 形状{shape_idx + 1} (索引{shape_idx}):")
            print(f"  形状类型: {shape.shape_type}")
            print(f"  形状名称: {getattr(shape, 'name', 'N/A')}")
            print(f"  文本内容: '{text}'")
            
            # 检查是否有标记文本
            if "【标记:" in text:
                print(f"  ✅ 包含标记")
            else:
                print(f"  ❌ 不包含标记")
        
        # 检查第8页（索引为7）
        slide = prs.slides[7]
        print(f"\n第8页形状数量: {len(slide.shapes)}")
        
        # 检查所有包含文本的形状
        text_shapes = []
        for i, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame is not None:
                all_text = ""
                for para in shape.text_frame.paragraphs:
                    all_text += para.text
                
                if all_text.strip():  # 只检查有实际文本内容的形状
                    text_shapes.append((i, shape, all_text))
        
        print(f"第8页包含文本的形状数量: {len(text_shapes)}")
        
        # 检查每个有文本的形状
        for shape_idx, shape, text in text_shapes:
            print(f"\n形状{shape_idx + 1} (索引{shape_idx}):")
            print(f"  形状类型: {shape.shape_type}")
            print(f"  形状名称: {getattr(shape, 'name', 'N/A')}")
            print(f"  文本内容: '{text}'")
            
            # 检查是否有标记文本
            if "【标记:" in text:
                print(f"  ✅ 包含标记")
                
                # 检查标记的样式
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "【标记:" in run.text:
                            print(f"    标记文本: '{run.text}'")
                            if run.font is not None:
                                print(f"    标记字体大小: {run.font.size}")
                                print(f"    标记字体颜色: {run.font.color.rgb if hasattr(run.font.color, 'rgb') else 'N/A'}")
                                print(f"    标记是否下划线: {run.font.underline}")
            else:
                print(f"  ❌ 不包含标记")
        
        # 检查首页是否有问题汇总
        first_slide = prs.slides[0]
        print(f"\n首页形状数量: {len(first_slide.shapes)}")
        
        # 查找问题汇总文本框
        summary_found = False
        for i, shape in enumerate(first_slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame is not None:
                text = shape.text_frame.text
                if "问题汇总" in text:
                    print(f"  ✅ 找到问题汇总 (形状{i}): {text}")
                    summary_found = True
                    break
        
        if not summary_found:
            print("  ❌ 未找到问题汇总")
            
    except Exception as e:
        print(f"检查失败: {e}")

if __name__ == "__main__":
    check_markings()
